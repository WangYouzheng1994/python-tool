#!/user/bin/env python
# -*-coding:utf-8-*-
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

# 设置字体
try:
    # 尝试指定可用字体
    plt.rcParams["font.family"] = ["SimHei", "WenQuanYi Micro Hei", "Heiti TC", "Arial"]
except:
    # 备选方案
    plt.rcParams["font.family"] = ["sans-serif"]
    plt.rcParams["axes.unicode_minus"] = False



# 读取Excel数据
data = pd.read_excel('input.xlsx')

# 假设数据有两列：产品和销售额
products = data['产品']
sales = data['销售额']

# 生成柱状图
plt.bar(products, sales)
plt.xlabel('产品')
plt.ylabel('销售额')
plt.title('产品销售额统计')

# 将图表保存到内存中
img_data = BytesIO()
plt.savefig(img_data, format='png')
img_data.seek(0)
plt.close()

# 创建PPT对象
prs = Presentation()

# 添加幻灯片布局
title_slide_layout = prs.slide_layouts[5]

# 添加幻灯片
slide = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide)

# 在幻灯片中添加图表图片
left = Inches(1)
top = Inches(1)
width = Inches(6)
height = Inches(4)
pic = slide.shapes.add_picture(img_data, left, top, width, height)

# 保存PPT
prs.save('report.pptx')